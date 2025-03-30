# SPDX-License-Identifier: GPL-3.0-or-later
# Copyright (C) 2025 Anant Patankar

"""
Document Processor for RAG Knowledge System

This module handles the transformation of raw documents into structured data
that can be efficiently processed by the embedding and retrieval stages.

As I've worked with various document types, I've found that preserving document
structure is critical for effective RAG systems. Simply splitting text into
chunks loses important context. This processor is designed to:

1. Extract hierarchical structure (sections, subsections) to maintain document organization
2. Identify and separately handle special content (tables, images, formulas, code)
3. Preserve metadata that enriches the semantic understanding
4. Clean up formatting issues that would reduce retrieval quality

Implementation notes:
- PDF processing uses PyMuPDF which I've found more reliable than alternatives
- Table detection uses geometry-based approaches rather than just text patterns
- The section detection algorithm uses both explicit (TOC) and implicit (formatting) cues
- For HTML/Markdown, we preserve the original structure rather than flattening

Performance considerations:
- Image extraction is the most memory-intensive operation
- For large documents (>1000 pages), consider processing in batches
- OCR is optional but dramatically improves searchability of image-heavy documents

Known limitations:
- Complex table structures in PDFs may not be perfectly preserved
- Heavily nested HTML can create section hierarchy issues
- Password-protected documents must be unlocked before processing

Author: Anant Patankar
Last updated: 30-March-2025
"""

import os
import re
import io
import json
import hashlib
import logging
import numpy as np
import pandas as pd
from typing import Dict, List, Any, Optional, Tuple, Union, Callable
from dataclasses import dataclass, field, asdict
from enum import Enum
from datetime import datetime
import time
import uuid
import traceback
from concurrent.futures import ThreadPoolExecutor, as_completed

import fitz  # PyMuPDF
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup
import markdown
import camelot  # For PDF tables
import docx  # For Word documents
import csv
import xml.etree.ElementTree as ET
import html2text
import pptx  # For PowerPoint presentations

import nltk
from nltk.tokenize import sent_tokenize, word_tokenize
from nltk.corpus import stopwords
import spacy
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity

# Embedding and vector libraries
from sentence_transformers import SentenceTransformer, util
import torch
import faiss
import hnswlib

from PIL import Image

from sklearn.metrics import precision_recall_fscore_support, ndcg_score

import sys
import os

from src.utils.helpers import generate_hash, ensure_directory, get_file_extension, \
                            clean_text

# Configure logging
logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(name)s - %(levelname)s - %(message)s')
logger = logging.getLogger('rag_system')

# Initialize required NLTK data
try:
    nltk.data.find('tokenizers/punkt')
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('punkt')
    nltk.download('stopwords')

# Initialize spaCy
try:
    nlp = spacy.load("en_core_web_sm")
except:
    logger.warning("spaCy model not found. Download with: python -m spacy download en_core_web_sm")
    # Fallback to simple NLP
    nlp = None

from src.config.config import ProcessingConfig
from src.config.doc_process_config import *

class DocumentProcessor:
    """
    My document processing engine that extracts structured content from various file formats.

    After working with several document parsing libraries, I created this processor to
    overcome limitations in how most RAG systems handle complex documents. Most systems
    treat documents as flat text, but I've found that preserving structural elements
    (sections, hierarchies, special content) dramatically improves retrieval quality.

    The processor handles multiple file formats with format-specific extraction logic,
    while maintaining a consistent output structure. I've optimized it to balance
    completeness of extraction with processing speed, with special attention to
    academic content like formulas, tables, and citations.

    Since different RAG applications have different needs, I've made the processor
    highly configurable through the ProcessingConfig class, allowing selective
    extraction of content types based on the use case.
    """

    def __init__(self, config: ProcessingConfig):
        """
        Initialize the document processor.

        Args:
            config: Processing configuration
        """
        self.config = config

        # Initialize document type handlers
        self.file_handlers = {
            '.pdf': self._process_pdf,
            '.epub': self._process_epub,
            '.txt': self._process_text,
            '.md': self._process_markdown,
            '.markdown': self._process_markdown,
            '.docx': self._process_docx,
            '.doc': self._process_docx,  # Will try to use docx handler
            '.html': self._process_html,
            '.htm': self._process_html,
            '.xml': self._process_xml,
            '.json': self._process_json,
            '.csv': self._process_csv,
            '.pptx': self._process_pptx,
            '.ppt': self._process_pptx,  # Will try to use pptx handler
        }

        # Custom Content Extractor Registration
        # ====================================

        # This section allows you to extend the document processor with your own file handlers.
        # You can add support for new file formats or override the built-in handlers for existing formats.

        # To use custom handlers:

        # 1. Create a handler function that follows this signature:
        #    def my_custom_handler(file_path: str, document: ProcessedDocument) -> None:
        #        """Process your custom format and populate the document object."""
        #        # Your processing code here

        # 2. Add your handler to the ProcessingConfig when initializing:
        #    config = ProcessingConfig(
        #        extract_tables=True,
        #        custom_content_extractors={
        #            '.xyz': my_custom_handler,  # Add handler for new format
        #            '.pdf': my_better_pdf_handler  # Override existing handler
        #        }
        #    )

        # 3. Initialize the DocumentProcessor with this config:
        #    processor = DocumentProcessor(config)

        # Your custom handlers should:
        # - Read the file at file_path
        # - Extract content, structure, and special elements
        # - Populate the document object with sections, tables, etc.
        # - Handle errors gracefully to prevent processing failures

        # Note: If you provide a handler for an extension that already has a built-in
        # handler (like .pdf), your custom handler will replace the original one.

        # The following loop registers all custom handlers from your configuration:
        for file_ext, handler in config.custom_content_extractors.items():
            self.file_handlers[file_ext] = handler

    def process_document(self, file_path: str) -> ProcessedDocument:
        """
        My document processing engine that extracts structured content from various file formats.

        After working with several document parsing libraries, I created this processor to
        overcome limitations in how most RAG systems handle complex documents. Most systems
        treat documents as flat text, but I've found that preserving structural elements
        (sections, hierarchies, special content) dramatically improves retrieval quality.

        The processor handles multiple file formats with format-specific extraction logic,
        while maintaining a consistent output structure. I've optimized it to balance
        completeness of extraction with processing speed, with special attention to
        academic content like formulas, tables, and citations.

        Since different RAG applications have different needs, I've made the processor
        highly configurable through the ProcessingConfig class, allowing selective
        extraction of content types based on the use case.
        """
        logger.info(f"Processing document: {file_path}")

        file_ext = get_file_extension(file_path)
        document = ProcessedDocument(source=file_path)

        # basic metadata addition
        document.metadata['filename'] = os.path.basename(file_path)
        document.metadata['file_type'] = file_ext.replace('.', '')
        document.metadata['file_size'] = os.path.getsize(file_path)
        document.metadata['processed_time'] = datetime.now().isoformat()

        try:
            # Finding appropriate handler for file type (File handler should be present in file handler)
            handler = self.file_handlers.get(file_ext)
            if handler:
                handler(file_path, document)
            else:
                # if no specific file handler present for the extension file will be processed as text
                logger.warning(f"No specific handler for {file_ext}, trying as text")
                self._process_text(file_path, document)

            # Clean and post-process content
            self._post_process(document)

            logger.info(f"Document processing complete: {len(document.content)} sections, "
                    f"{len(document.tables)} tables, {len(document.images)} images, "
                    f"{len(document.formulas)} formulas, {len(document.code_blocks)} code blocks")

            return document

        except Exception as e:
            logger.error(f"Error processing document {file_path}: {e}", exc_info=True)
            # if data is partially process that will returned if available
            return document

    def _process_pdf(self, file_path: str, processed_document: ProcessedDocument) -> None:
        """
            Extract content from PDF with section hierarchy preservation.

            Uses PyMuPDF to extract text, tables, images, and document structure
            for optimal retrieval context.
        """
        pdf_doc = fitz.open(file_path)

        # Extract metadata
        metadata = pdf_doc.metadata
        if metadata:
            processed_document.metadata.update({
                'title': metadata.get('title', ''),
                'author': metadata.get('author', ''),
                'subject': metadata.get('subject', ''),
                'keywords': metadata.get('keywords', ''),
                'page_count': len(pdf_doc)
            })

        # Initialize sections_by_page as an empty dict by default,
        # this is added because I was getting unboundlocal error
        sections_by_page = {}

        # Extract TOC/outline if available (TOC: Table Of Content)
        toc = pdf_doc.get_toc()
        if toc:
            structure = [{'level': t[0], 'title': t[1], 'page': t[2]} for t in toc]
            processed_document.metadata['toc'] = structure

            # Use TOC to create initial section structure
            current_section = None
            sections_by_page = {}

            for item in structure:
                section = Section(
                    title=item['title'],
                    content="",
                    level=item['level'],
                    page_start=item['page'] - 1  # PDF TOC is 1-indexed
                )

                # Record section by page for content assignment
                if item['page'] not in sections_by_page:
                    sections_by_page[item['page'] - 1] = []
                sections_by_page[item['page'] - 1].append(section)

                processed_document.content.append(section)

            # Set end pages for sections
            for i in range(len(processed_document.content) - 1):
                processed_document.content[i].page_end = processed_document.content[i+1].page_start - 1

            # Set end page for last section
            if processed_document.content:
                processed_document.content[-1].page_end = len(pdf_doc) - 1

        # page by page content extraction
        for page_num, page in enumerate(pdf_doc):
            # Get text content
            text = page.get_text()

            # Clean headers and footers if enabled
            if self.config.clean_header_footer:
                text = self._clean_header_footer(text, page_num, len(pdf_doc))

            # Assign text to appropriate section if we have TOC
            if sections_by_page and page_num in sections_by_page:
                for section in sections_by_page[page_num]:
                    section.content += text
            elif sections_by_page:
                # Find the section this page belongs to
                for section in reversed(processed_document.content):
                    if section.page_start <= page_num and (section.page_end is None or section.page_end >= page_num):
                        section.content += text
                        break
            else:
                # No TOC, use heuristics to detect sections
                lines = text.split('\n')
                if lines and page_num == 0:
                    # First page might have title
                    title = lines[0].strip()
                    content = '\n'.join(lines[1:])
                    processed_document.content.append(Section(
                        title=title, 
                        content=content,
                        page_start=page_num,
                        page_end=page_num
                    ))
                elif processed_document.content:
                    # Append to last section
                    processed_document.content[-1].content += '\n' + text
                    processed_document.content[-1].page_end = page_num
                else:
                    # Create a new section
                    processed_document.content.append(Section(
                        title=f"Page {page_num+1}",
                        content=text,
                        page_start=page_num,
                        page_end=page_num
                    ))

            # Extract tables if requested
            if self.config.extract_tables:
                self._extract_tables_from_pdf_page(file_path, page_num, processed_document)

            # Extract images if requested
            if self.config.extract_images:
                self._extract_images_from_pdf_page(pdf_doc, page, page_num, processed_document)

            # Extract links
            if self.config.extract_links:
                self._extract_links_from_pdf_page(page, page_num, processed_document)

        # Extract formulas if requested
        if self.config.extract_formulas:
            self._extract_formulas_from_document(processed_document)

        # Extract code blocks if requested
        if self.config.extract_code_blocks:
            self._extract_code_blocks_from_document(processed_document)

        # Extract references if requested
        if self.config.extract_references:
            self._extract_references_from_document(processed_document)

        pdf_doc.close()

    def _extract_tables_from_pdf_page(self, file_path: str, page_num: int, document: ProcessedDocument) -> None:
        """
        Extract tables from PDF page using Camelot library.

        Identifies tabular data and converts to structured format with
        headers and content rows for improved searchability.
        """
        try:
            if self.config.table_extraction_method == "automatic":
                # Try using camelot module for table extraction
                tables = camelot.read_pdf(file_path, pages=str(page_num + 1), flavor='lattice')

                for i, table in enumerate(tables):
                    if table.df.empty:
                        continue

                    # Convert to table format
                    table_data = TableData(
                        content=table.df.values.tolist(),
                        headers=table.df.columns.tolist(),
                        page_num=page_num,
                        position=i,
                        caption=f"Table on page {page_num+1}"
                    )

                    document.tables.append(table_data)
            elif self.config.table_extraction_method == "heuristic":
                # Alternative basic heuristic table detection
                # This is simplified and would need to be expanded 
                # this should be expanded for real use
                for section in document.content:
                    if section.page_start <= page_num <= (section.page_end or page_num):
                        text = section.content

                        # Very simple heuristic - look for multiple pipe or tab characters
                        table_candidates = re.findall(r'([^\n]+\|[^\n]+\|[^\n]+)', text)
                        for i, candidate in enumerate(table_candidates):
                            rows = [line.strip() for line in candidate.split('\n') if line.strip()]
                            if len(rows) > 1:  # Need at least header and one data row
                                # Split by pipe or tab
                                delimiter = '|' if '|' in rows[0] else '\t'
                                content = [row.split(delimiter) for row in rows]
                                if content and all(len(row) == len(content[0]) for row in content):
                                    table_data = TableData(
                                        content=content[1:],
                                        headers=content[0] if len(content) > 1 else [],
                                        page_num=page_num,
                                        position=i,
                                        section_id=section.section_id
                                    )
                                    document.tables.append(table_data)
        except Exception as e:
            logger.warning(f"Error extracting tables from page {page_num+1}: {e}")

    def _extract_images_from_pdf_page(self, doc, page, page_num: int, document: ProcessedDocument) -> None:
        """
            Image extraction from a PDF page. This method captures image data, dimensions, and 
            OCR can be ENABLED to extract image conent for searchable image content
        """
        try:
            image_list = page.get_images(full=True)
            for img_index, img in enumerate(image_list):
                xref = img[0]
                base_image = doc.extract_image(xref)

                if base_image and len(base_image["image"]) <= self.config.max_image_size:
                    # Check for nearby caption (simple heuristic approach)
                    caption = ""
                    try:
                        text = page.get_text()
                        caption_pattern = r"(Figure|Fig\.?)\s+\d+[.:]\s*([^\n]+)"
                        caption_matches = re.findall(caption_pattern, text)
                        if caption_matches:
                            caption = caption_matches[0][1].strip()
                    except:
                        pass

                    # Perform OCR if enabled
                    ocr_text = ""
                    if self.config.ocr_enabled:
                        try:
                            from PIL import Image
                            import pytesseract

                            # Convert image data to PIL Image
                            image = Image.open(io.BytesIO(base_image["image"]))

                            # Perform OCR
                            ocr_text = pytesseract.image_to_string(
                                image, 
                                lang=self.config.ocr_language
                            )
                        except Exception as ocr_error:
                            logger.warning(f"OCR error: {ocr_error}")

                    image_data = ImageData(
                        data=base_image["image"],
                        extension=base_image["ext"],
                        page_num=page_num,
                        caption=caption,
                        ocr_text=ocr_text
                    )

                    document.images.append(image_data)
        except Exception as e:
            logger.warning(f"Error extracting images from page {page_num+1}: {e}")

    def _extract_links_from_pdf_page(self, page, page_num: int, document: ProcessedDocument) -> None:
        try:
            links = page.get_links()
            for link in links:
                if 'uri' in link:  # for external URL
                    # Get the text near the link (simple way)
                    text = "Link"  # Default text
                    rect = link.get('rect')
                    if rect:
                        # Get text from the area of the link
                        try:
                            text_area = page.get_text("text", clip=rect)
                            if text_area.strip():
                                text = text_area.strip()
                        except:
                            pass

                    link_data = LinkData(
                        text=text,
                        url=link['uri'],
                        page_num=page_num,
                        link_type="external"
                    )
                    document.links.append(link_data)
                elif 'page' in link:  # for internal link to another page
                    # Similar approach to get text
                    text = "Internal link"
                    rect = link.get('rect')
                    if rect:
                        try:
                            text_area = page.get_text("text", clip=rect)
                            if text_area.strip():
                                text = text_area.strip()
                        except:
                            pass

                    link_data = LinkData(
                        text=text,
                        url=f"#page={link['page']+1}",  # +1 because PDF pages are zero-indexed
                        page_num=page_num,
                        link_type="internal"
                    )
                    document.links.append(link_data)
        except Exception as e:
            logger.warning(f"Error extracting links from page {page_num+1}: {e}")

    def _process_epub(self, file_path: str, document: ProcessedDocument) -> None:
        """
            This method processes epub files.
        """
        book = epub.read_epub(file_path)

        # Extract metadata
        title = book.get_metadata('DC', 'title')
        creator = book.get_metadata('DC', 'creator')
        language = book.get_metadata('DC', 'language')

        document.metadata.update({
            'title': title[0][0] if title else '',
            'author': creator[0][0] if creator else '',
            'language': language[0][0] if language else '',
            'format': 'EPUB'
        })

        # Process content
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                # Getting content
                content = item.get_content().decode('utf-8')
                soup = BeautifulSoup(content, 'html.parser')

                # Extract title from heading elements
                title = ""
                heading = soup.find(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])
                if heading:
                    title = heading.get_text().strip()
                else:
                    # Try to find titile by other ways
                    title_elem = soup.find('title')
                    if title_elem:
                        title = title_elem.get_text().strip()
                    else:
                        title = f"Item {item.get_id()}"

                # Clean text content
                text = soup.get_text().strip()

                # Create section
                if text:
                    section = Section(
                        title=title,
                        content=text,
                        metadata={'item_id': item.get_id()}
                    )
                    document.content.append(section)

                # Process the rest of the HTML content (tables, images, etc.)
                self._process_html_content(soup, section, document, item, book)

    def _process_html_content(self, soup, section, document, item=None, book=None):
        """
            This function processes HTML content for
                - tables, 
                - images, 
                - formulas, 
                - code, etc.
            Note: Regular HTML processing it not implemented yet, but added 
                    placeholder and explanatory comment for it.
        """
        section_id = section.section_id if hasattr(section, 'section_id') else None

        # to extract tables if requested
        if self.config.extract_tables:
            for idx, table_elem in enumerate(soup.find_all('table')):
                # extract rows and cells
                rows = []
                headers = []

                # Try to get headers of table first 
                thead = table_elem.find('thead')
                if thead:
                    header_row = thead.find('tr')
                    if header_row:
                        headers = [th.get_text().strip() for th in header_row.find_all(['th', 'td'])]

                # If no headers in thead, use first row
                if not headers:
                    first_row = table_elem.find('tr')
                    if first_row and first_row.find('th'):
                        headers = [th.get_text().strip() for th in first_row.find_all('th')]

                # Get table rows
                for tr in table_elem.find_all('tr'):
                    row = [td.get_text().strip() for td in tr.find_all(['td', 'th'])]
                    if row:
                        rows.append(row)

                # Find caption if any
                caption = ""
                caption_elem = table_elem.find('caption')
                if caption_elem:
                    caption = caption_elem.get_text().strip()

                if rows:
                    table_data = TableData(
                        content=rows,
                        headers=headers,
                        section_id=section_id,
                        caption=caption
                    )
                    document.tables.append(table_data)

        # Extract images if requested
        if self.config.extract_images:
            for idx, img in enumerate(soup.find_all('img')):
                src = img.get('src', '')
                alt = img.get('alt', '')

                # Try to find the actual image data
                if src and book and item:  # EPUB specific
                    try:
                        image_item = book.get_item_with_href(src)
                        if image_item:
                            img_data = image_item.get_content()

                            # Determine extension from mime type
                            mime = image_item.get_type()
                            ext = mime.split('/')[-1] if '/' in mime else 'jpg'

                            # Find figure caption if any
                            caption = ""
                            fig_elem = img.find_parent('figure')
                            if fig_elem:
                                figcaption = fig_elem.find('figcaption')
                                if figcaption:
                                    caption = figcaption.get_text().strip()

                            image_data = ImageData(
                                data=img_data,
                                extension=ext,
                                caption=caption,
                                alt_text=alt,
                                section_id=section_id
                            )
                            document.images.append(image_data)
                    except Exception as e:
                        logger.warning(f"Error extracting image {src}: {e}")
                elif src:  # Regular HTML
                    # This is added for future processing implementation, see comments
                    # For regular HTML it is needed to fetch the image, but this is skipped that here
                    # This would involve downloading the image or handling local references
                    # Further processing will be added here
                    pass

        # Extract formulas if requested
        if self.config.extract_formulas:
            # check for MathML
            for idx, math_elem in enumerate(soup.find_all('math')):
                formula_data = FormulaData(
                    content=str(math_elem),
                    is_inline=math_elem.get('display') != 'block',
                    section_id=section_id
                )
                document.formulas.append(formula_data)

            # checnking for TeX formulas
            text_content = soup.get_text()
            formula_pattern = r'(?:\${1,2})(.*?)(?:\${1,2})'
            for match in re.finditer(formula_pattern, text_content, re.DOTALL):
                formula_text = match.group(1)
                is_inline = not (match.group(0).startswith('$') and match.group(0).endswith('$'))

                formula_data = FormulaData(
                    content=formula_text,
                    is_inline=is_inline,
                    position=match.start(),
                    section_id=section_id
                )
                document.formulas.append(formula_data)

        # Extract code blocks if requested
        if self.config.extract_code_blocks:
            for idx, code_elem in enumerate(soup.find_all(['code', 'pre'])):
                # Try to determine the language
                language = ""
                if 'class' in code_elem.attrs:
                    classes = code_elem['class']
                    for cls in classes:
                        if cls.startswith('language-') or cls.startswith('lang-'):
                            language = cls.split('-', 1)[1]
                            break

                # Get the code content
                code_content = code_elem.get_text()

                # Skip if too short or empty
                if not code_content or len(code_content) < 10:
                    continue

                # Find caption if any
                caption = ""
                fig_elem = code_elem.find_parent('figure')
                if fig_elem:
                    figcaption = fig_elem.find('figcaption')
                    if figcaption:
                        caption = figcaption.get_text().strip()

                code_data = CodeBlockData(
                    content=code_content,
                    language=language,
                    section_id=section_id,
                    caption=caption
                )
                document.code_blocks.append(code_data)

        # Link extraction if requested
        if self.config.extract_links:
            for idx, link_elem in enumerate(soup.find_all('a')):
                href = link_elem.get('href', '')
                if href:
                    link_text = link_elem.get_text().strip()
                    if not link_text:
                        link_text = href

                    # Determine link type (external, internal, anchor link)
                    link_type = "external"
                    if href.startswith('#'):
                        link_type = "anchor"
                    elif not (href.startswith('http://') or href.startswith('https://')):
                        link_type = "internal"

                    link_data = LinkData(
                        text=link_text,
                        url=href,
                        section_id=section_id,
                        link_type=link_type
                    )
                    document.links.append(link_data)

    def _process_text(self, file_path: str, document: ProcessedDocument) -> None:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
            text = file.read()

        # Basic metadata
        document.metadata.update({
            'title': os.path.basename(file_path),
            'format': 'Plain Text',
            'character_count': len(text)
        })

        # try to detect sections based on line patterns
        lines = text.split('\n')
        section_texts = []
        current_section = {"title": document.metadata['title'], "content": ""}

        # simple and heuristics section detection
        for i, line in enumerate(lines):
            line_strip = line.strip()

            # Skip empty lines
            if not line_strip:
                current_section["content"] += "\n"
                continue

            # Potential section header heuristics:
            # 1. All caps line under certain length
            # 2. Line followed by underline chars (=, -, ~, etc.)
            # 3. Numbered headings (1., 1.1., Chapter 1:, etc.)
            # 4. Short line followed by blank line
            is_heading = False

            # Check for all caps with reasonable length
            if line_strip.isupper() and 3 < len(line_strip) < 100:
                is_heading = True

            # Check for heading followed by separator line
            elif i < len(lines) - 1 and lines[i+1].strip() and all(c == lines[i+1].strip()[0] for c in lines[i+1].strip()):
                is_heading = True

            # Check for numbered headings
            elif re.match(r'^(\d+\.)+\s+\S.*$|^(Chapter|Section|Part)\s+\d+', line_strip):
                is_heading = True

            # Check for short line followed by blank line
            elif len(line_strip) < 50 and i < len(lines) - 1 and not lines[i+1].strip():
                is_heading = True

            if is_heading:
                # Save previous section if it has content
                if current_section["content"].strip():
                    section_texts.append(current_section.copy())

                # start new section
                current_section = {"title": line_strip, "content": ""}
            else:
                current_section["content"] += line + "\n"

        # Add the last section
        if current_section["content"].strip():
            section_texts.append(current_section)

        # if Unable to detect sections, whole text will be used
        if not section_texts:
            section_texts = [{
                "title": document.metadata['title'],
                "content": text
            }]

        # Create section objects
        for section_dict in section_texts:
            section = Section(
                title=section_dict["title"],
                content=section_dict["content"]
            )
            document.content.append(section)

        # Detect and extract other content
        if self.config.extract_tables:
            self._extract_tables_from_text(text, document)

        if self.config.extract_code_blocks:
            self._extract_code_blocks_from_text(text, document)

        if self.config.extract_formulas:
            self._extract_formulas_from_document(document)

        if self.config.extract_references:
            self._extract_references_from_document(document)

        if self.config.extract_links:
            self._extract_links_from_text(text, document)

    def _process_markdown(self, file_path: str, document: ProcessedDocument) -> None:
        """
            Process Markdown file.
            Markdown file: file with .md extension
        """
        with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
            md_text = file.read()

        # convert to HTML for better processing
        html = markdown.markdown(md_text, extensions=['tables', 'fenced_code'])
        soup = BeautifulSoup(html, 'html.parser')

        # Extract basic metadata
        document.metadata.update({
            'title': os.path.basename(file_path),
            'format': 'Markdown',
            'character_count': len(md_text)
        })

        # Check for Title in first h1 if available
        title_elem = soup.find('h1')
        if title_elem:
            document.metadata['title'] = title_elem.get_text().strip()

        # Extract sections based on headers
        sections = []
        current_section = None
        current_content = []

        # Process each element
        for elem in soup.children:
            if hasattr(elem, 'name') and elem.name in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']:
                # Save previous section if exists
                if current_section:
                    level = int(current_section.name[1])
                    title = current_section.get_text().strip()
                    content = '\n'.join(current_content)
                    sections.append((level, title, content))

                # Start new section
                current_section = elem
                current_content = []
            else:
                if current_section:
                    current_content.append(str(elem))
                else:
                    # Content before first heading
                    current_content.append(str(elem))

        # Add the last section
        if current_section:
            level = int(current_section.name[1])
            title = current_section.get_text().strip()
            content = '\n'.join(current_content)
            sections.append((level, title, content))
        elif current_content:
            # IF No sections found, Use filename as title
            sections.append((1, document.metadata['title'], '\n'.join(current_content)))

        # Create section objects
        for level, title, content in sections:
            # Clean content - BeautifulSoup might give us some HTML
            clean_soup = BeautifulSoup(content, 'html.parser')
            clean_text = clean_soup.get_text()

            section = Section(
                title=title,
                content=clean_text,
                level=level
            )
            document.content.append(section)

        # Process all HTML content
        self._process_html_content(soup, None, document)

        # check for code blocks in original markdown
        if self.config.extract_code_blocks:
            code_block_pattern = r'```(\w*)\n(.*?)\n```'
            for match in re.finditer(code_block_pattern, md_text, re.DOTALL):
                language = match.group(1).strip()
                code = match.group(2).strip()

                if code:
                    code_data = CodeBlockData(
                        content=code,
                        language=language,
                        position=match.start()
                    )
                    document.code_blocks.append(code_data)

    def _process_docx(self, file_path: str, document: ProcessedDocument) -> None:
        try:
            doc = docx.Document(file_path)

            # extract metadata
            core_properties = doc.core_properties
            document.metadata.update({
                'title': core_properties.title or os.path.basename(file_path),
                'author': core_properties.author or '',
                'created': core_properties.created.isoformat() if core_properties.created else '',
                'modified': core_properties.modified.isoformat() if core_properties.modified else '',
                'last_modified_by': core_properties.last_modified_by or '',
                'revision': core_properties.revision or '',
                'format': 'DOCX'
            })

            # paragraph proccessing and building sections
            current_section = None
            sections = []

            for para in doc.paragraphs:
                text = para.text.strip()
                if not text:
                    continue

                # Check if paragraph is a heading
                if para.style.name.startswith('Heading'):
                    # Extract level of heading
                    level = 1
                    try:
                        level = int(para.style.name.replace('Heading', ''))
                    except ValueError:
                        pass

                    # Save previous section if exists
                    if current_section:
                        sections.append(current_section)

                    # Create new section
                    current_section = {
                        'title': text,
                        'content': '',
                        'level': level
                    }
                elif current_section:
                    # Add to current section content
                    current_section['content'] += text + '\n\n'
                else:
                    # No section yet, create default
                    current_section = {
                        'title': document.metadata['title'],
                        'content': text + '\n\n',
                        'level': 1
                    }

            # Add the last section
            if current_section:
                sections.append(current_section)

            # If no sections, create a single section with whole content
            if not sections:
                all_text = '\n\n'.join([p.text for p in doc.paragraphs if p.text.strip()])
                sections.append({
                    'title': document.metadata['title'],
                    'content': all_text,
                    'level': 1
                })

            # Create section objects
            for section_dict in sections:
                section = Section(
                    title=section_dict['title'],
                    content=section_dict['content'].strip(),
                    level=section_dict['level']
                )
                document.content.append(section)

            # Extract tables if requested
            if self.config.extract_tables:
                for i, table in enumerate(doc.tables):
                    rows = []
                    headers = []

                    # Assuming first row might be header
                    if table.rows:
                        headers = [cell.text.strip() for cell in table.rows[0].cells]

                    # Get content rows
                    for row in table.rows[1:]:
                        row_data = [cell.text.strip() for cell in row.cells]
                        if any(row_data):  # skipping empty rows
                            rows.append(row_data)

                    if rows:
                        table_data = TableData(
                            content=rows,
                            headers=headers,
                            position=i,
                            caption=f"Table {i+1}"
                        )
                        document.tables.append(table_data)

            # Text processing for formulas, code blocks, references, etc.
            all_text = '\n\n'.join([p.text for p in doc.paragraphs if p.text.strip()])
            if self.config.extract_formulas:
                self._extract_formulas_from_text(all_text, document)

            if self.config.extract_code_blocks:
                self._extract_code_blocks_from_text(all_text, document)

            if self.config.extract_references:
                self._extract_references_from_text(all_text, document)

            if self.config.extract_links:
                # Extract hyperlinks from document
                for para in doc.paragraphs:
                    for run in para.runs:
                        if run.hyperlink:
                            link_data = LinkData(
                                text=run.text,
                                url=run.hyperlink.url,
                                link_type="external"
                            )
                            document.links.append(link_data)

        except Exception as e:
            logger.error(f"Error processing DOCX file {file_path}: {e}")

    def _process_html(self, file_path: str, document: ProcessedDocument) -> None:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
            html_content = file.read()

        soup = BeautifulSoup(html_content, 'html.parser')

        # Extract basic metadata
        title_elem = soup.find('title')
        document.metadata.update({
            'title': title_elem.get_text().strip() if title_elem else os.path.basename(file_path),
            'format': 'HTML',
            'character_count': len(html_content)
        })

        # Meta tags extraction
        for meta in soup.find_all('meta'):
            name = meta.get('name') or meta.get('property')
            content = meta.get('content')
            if name and content:
                document.metadata[name] = content

        # Extract sections based on headers
        headers = soup.find_all(['h1', 'h2', 'h3', 'h4', 'h5', 'h6'])

        if headers:
            # Process content based on headers
            for i, header in enumerate(headers):
                # Get header information
                title = header.get_text().strip()
                level = int(header.name[1])

                # Get content until we reach to next header
                content = []
                elem = header.next_sibling
                while elem and (not hasattr(elem, 'name') or elem.name not in ['h1', 'h2', 'h3', 'h4', 'h5', 'h6']):
                    if hasattr(elem, 'get_text'):
                        content.append(elem.get_text())
                    elif isinstance(elem, str):
                        content.append(elem)
                    elem = elem.next_sibling

                # Create section
                if content:
                    section = Section(
                        title=title,
                        content='\n'.join(content).strip(),
                        level=level
                    )
                    document.content.append(section)
        else:
            # No headers found, extract text from body
            body = soup.find('body')
            if body:
                # For simple pages, just get all text
                text = body.get_text().strip()
                section = Section(
                    title=document.metadata['title'],
                    content=text,
                    level=1
                )
                document.content.append(section)

        # Process HTML content for tables, images, links, etc.
        self._process_html_content(soup, None, document)

    def _process_csv(self, file_path: str, document: ProcessedDocument) -> None:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
            reader = csv.reader(file)
            rows = list(reader)

        if not rows:
            return

        # Extract basic metadata
        document.metadata.update({
            'title': os.path.basename(file_path),
            'format': 'CSV',
            'row_count': len(rows),
            'column_count': len(rows[0]) if rows else 0
        })

        # A single section with basic info will be created
        summary = f"CSV data with {len(rows)} rows and {len(rows[0]) if rows else 0} columns."
        section = Section(
            title=document.metadata['title'],
            content=summary
        )
        document.content.append(section)

        # Creating table data
        headers = rows[0] if len(rows) > 0 else []
        content = rows[1:] if len(rows) > 1 else []

        table_data = TableData(
            content=content,
            headers=headers,
            caption=f"Data from {os.path.basename(file_path)}"
        )
        document.tables.append(table_data)

    def _process_json(self, file_path: str, document: ProcessedDocument) -> None:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
            try:
                data = json.load(file)
            except json.JSONDecodeError:
                logger.error(f"Invalid JSON in {file_path}")
                return

        # Extract basic metadata
        document.metadata.update({
            'title': os.path.basename(file_path),
            'format': 'JSON'
        })

        # Create a single section with the JSON content
        if isinstance(data, dict):
            # For dictionaries, list keys at the top level
            keys = list(data.keys())
            content = f"JSON object with {len(keys)} keys: {', '.join(keys[:10])}"
            if len(keys) > 10:
                content += f" and {len(keys) - 10} more"

            section = Section(
                title=document.metadata['title'],
                content=content
            )
            document.content.append(section)

            # Extract arrays as tables if possible
            for key, value in data.items():
                if isinstance(value, list) and value and isinstance(value[0], dict):
                    self._extract_table_from_json_array(value, key, document)

        elif isinstance(data, list):
            # For arrays, describe the length and type
            content = f"JSON array with {len(data)} items"
            if data and isinstance(data[0], dict):
                # Try to extract as table if it's an array of objects
                self._extract_table_from_json_array(data, "Data", document)

            section = Section(
                title=document.metadata['title'],
                content=content
            )
            document.content.append(section)

        # Add the prettified JSON as a code block
        if self.config.extract_code_blocks:
            code_data = CodeBlockData(
                content=json.dumps(data, indent=2),
                language="json",
                caption=f"JSON content from {os.path.basename(file_path)}"
            )
            document.code_blocks.append(code_data)

    def _extract_table_from_json_array(self, data: List[Dict], name: str, document: ProcessedDocument) -> None:
        if not data or not isinstance(data[0], dict):
            return

        # Collecting all possible keys from json
        all_keys = set()
        for item in data:
            all_keys.update(item.keys())

        # Create headers from keys
        headers = list(all_keys)

        # Create content rows
        content = []
        for item in data:
            row = [str(item.get(key, '')) for key in headers]
            content.append(row)

        # Create table data
        if content:
            table_data = TableData(
                content=content,
                headers=headers,
                caption=f"{name} from JSON"
            )
            document.tables.append(table_data)

    def _process_xml(self, file_path: str, document: ProcessedDocument) -> None:
        try:
            tree = ET.parse(file_path)
            root = tree.getroot()

            # Extract basic metadata
            document.metadata.update({
                'title': os.path.basename(file_path),
                'format': 'XML',
                'root_element': root.tag
            })

            # Create a section with basic XML structure info
            namespaces = []
            if root.attrib:
                for key, value in root.attrib.items():
                    if key.startswith('xmlns:'):
                        namespaces.append(f"{key}=\"{value}\"")

            content = f"XML document with root element <{root.tag}>"
            if namespaces:
                content += f"\nNamespaces: {', '.join(namespaces)}"

            # Count children
            direct_children = list(root)
            child_tags = {}
            for child in direct_children:
                child_tags[child.tag] = child_tags.get(child.tag, 0) + 1

            content += f"\nDirect children of root: {len(direct_children)}"
            content += "\nChild elements:"
            for tag, count in child_tags.items():
                content += f"\n- {tag}: {count} elements"

            section = Section(
                title=document.metadata['title'],
                content=content
            )
            document.content.append(section)

            # Extract repeating elements as tables if possible
            if self.config.extract_tables:
                self._extract_tables_from_xml(root, document)

            # Add prettified XML as code block
            if self.config.extract_code_blocks:
                xml_str = ET.tostring(root, encoding='unicode', method='xml')
                # Basic indentation for prettier display
                level = 0
                pretty_xml = []
                for line in xml_str.split('>'):
                    if not line.strip():
                        continue
                    if line.find('</') >= 0:
                        level -= 1
                    pretty_xml.append('  ' * level + line + '>')
                    if line.find('/>') < 0 and line.find('</') < 0:
                        level += 1

                code_data = CodeBlockData(
                    content='\n'.join(pretty_xml),
                    language="xml",
                    caption=f"XML content from {os.path.basename(file_path)}"
                )
                document.code_blocks.append(code_data)

        except Exception as e:
            logger.error(f"Error processing XML file {file_path}: {e}")

    def _extract_tables_from_xml(self, element: ET.Element, document: ProcessedDocument) -> None:
        """
            This function processes XML file to extract tables from repeating XML elements.
        """
        # Find child elements that repeat and have a consistent structure
        child_elements = {}
        for child in element:
            if child.tag not in child_elements:
                child_elements[child.tag] = []
            child_elements[child.tag].append(child)

        # For each repeating element type, check if they could form a table
        for tag, elements in child_elements.items():
            if len(elements) <= 1:
                continue

            # Check if all elements have similar structure - get all possible sub-element names
            all_keys = set()
            for element in elements:
                # Get direct child tags
                for child in element:
                    all_keys.add(child.tag)
                # Also add attributes
                for attr in element.attrib:
                    all_keys.add(f"@{attr}")

            # If we have a reasonable number of keys, create a table
            if all_keys and len(all_keys) <= 50:  # Arbitrary limit to avoid huge tables
                headers = list(all_keys)

                # Create content rows
                content = []
                for element in elements:
                    row = []
                    for key in headers:
                        if key.startswith('@'):  # Attribute
                            attr_name = key[1:]
                            row.append(element.attrib.get(attr_name, ''))
                        else:  # Child element
                            # Find matching child and get its text
                            child = element.find(key)
                            row.append(child.text.strip() if child is not None and child.text else '')
                    content.append(row)

                # Create table data
                table_data = TableData(
                    content=content,
                    headers=headers,
                    caption=f"Data from {tag} elements"
                )
                document.tables.append(table_data)

    def _process_pptx(self, file_path: str, document: ProcessedDocument) -> None:
        try:
            presentation = pptx.Presentation(file_path)

            # Extract basic metadata
            document.metadata.update({
                'title': os.path.basename(file_path),
                'format': 'PPTX',
                'slide_count': len(presentation.slides)
            })

            # Process slides
            for i, slide in enumerate(presentation.slides):
                # Extract slide title
                title = ""
                if slide.shapes.title:
                    title = slide.shapes.title.text.strip()
                if not title:
                    title = f"Slide {i+1}"

                # Extract text content from all shapes
                texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            texts.append(shape.text.strip())

                content = "\n\n".join(texts)

                # Create section for this slide
                section = Section(
                    title=title,
                    content=content,
                    level=1,  # All slides at same level
                    metadata={'slide_number': i+1}
                )
                document.content.append(section)

                # Extract tables if requested
                if self.config.extract_tables:
                    for shape in slide.shapes:
                        if shape.has_table:
                            table = shape.table
                            rows = []
                            for row in table.rows:
                                row_data = [cell.text.strip() for cell in row.cells]
                                rows.append(row_data)

                            # Create table data
                            if rows:
                                # Use first row as header if it looks different
                                headers = []
                                content = rows

                                table_data = TableData(
                                    content=content,
                                    headers=headers,
                                    section_id=section.section_id,
                                    caption=f"Table on slide {i+1}"
                                )
                                document.tables.append(table_data)

                # Extract images if requested
                if self.config.extract_images:
                    for shape in slide.shapes:
                        if hasattr(shape, 'image'):
                            try:
                                # Get image data
                                image = shape.image
                                image_bytes = image.blob

                                # Determine extension from content type
                                content_type = image.content_type
                                ext = content_type.split('/')[-1] if content_type else 'jpg'

                                image_data = ImageData(
                                    data=image_bytes,
                                    extension=ext,
                                    section_id=section.section_id,
                                    caption=f"Image on slide {i+1}"
                                )
                                document.images.append(image_data)
                            except Exception as img_err:
                                logger.warning(f"Error extracting image on slide {i+1}: {img_err}")

        except Exception as e:
            logger.error(f"Error processing PPTX file {file_path}: {e}")

    def _extract_tables_from_text(self, text: str, document: ProcessedDocument) -> None:
        """
            If normal text (.txt) file contains tables this function will be used.
        """
        # Look for common table patterns in plain text
        # ASCII tables with bars or plus signs
        table_pattern = r'[|\+][-+]+[|\+]'
        for match in re.finditer(table_pattern, text, re.MULTILINE):
            # Extract the full table - go back a few lines and forward until pattern breaks
            start_pos = match.start()
            while start_pos > 0 and text[start_pos] != '\n':
                start_pos -= 1

            end_pos = match.end()
            while end_pos < len(text) and text[end_pos] != '\n':
                end_pos += 1

            # Get 10 lines around the match to capture the full table
            lines_before = text[:start_pos].split('\n')[-5:]
            table_line = text[start_pos:end_pos]
            lines_after = text[end_pos:].split('\n')[:5]

            table_text = '\n'.join(lines_before + [table_line] + lines_after)

            # Parse table - split by lines and then by pipe or plus
            rows = []
            for line in table_text.split('\n'):
                if '|' in line and not all(c in '+-|' for c in line.strip()):
                    cells = [cell.strip() for cell in line.split('|')]
                    # Remove empty cells at start/end from split
                    if cells and not cells[0]:
                        cells = cells[1:]
                    if cells and not cells[-1]:
                        cells = cells[:-1]
                    if cells:
                        rows.append(cells)

            if rows:
                # First row might be headers
                headers = rows[0] if len(rows) > 1 else []

                table_data = TableData(
                    content=rows[1:] if len(rows) > 1 else rows,
                    headers=headers,
                    position=match.start()
                )
                document.tables.append(table_data)

    def _extract_code_blocks_from_text(self, text: str, document: ProcessedDocument) -> None:
        """Extract code blocks from text content."""
        if self.config.code_extraction_method == "pattern":
            # check if indented blocks present
            lines = text.split('\n')
            in_code_block = False
            current_block = []

            for line in lines:
                # If line starts with 4+ spaces or a tab, it might be code
                if re.match(r'^(\t|    )', line):
                    if not in_code_block:
                        in_code_block = True
                    current_block.append(line)
                else:
                    # End of indented block
                    if in_code_block and current_block:
                        code = '\n'.join(current_block)
                        if len(code.strip()) > 30:  # Ignore tiny blocks
                            code_data = CodeBlockData(
                                content=code,
                                language=""  # Unknown language
                            )
                            document.code_blocks.append(code_data)
                        current_block = []
                        in_code_block = False

            # Check for code blocks in backticks like ```python code ```
            code_block_pattern = r'```(\w*)\n(.*?)\n```'
            for match in re.finditer(code_block_pattern, text, re.DOTALL):
                language = match.group(1).strip()
                code = match.group(2).strip()

                if code:
                    code_data = CodeBlockData(
                        content=code,
                        language=language,
                        position=match.start()
                    )
                    document.code_blocks.append(code_data)

    def _extract_formulas_from_document(self, document: ProcessedDocument) -> None:
        """
            Mathematical formula extraction from document content
        """
        if self.config.formula_extraction_method == "pattern":
            # Iterate through sections to find formulas
            for section in document.content:
                # Look for LaTeX style formulas ($...$, $...$)
                formula_pattern = r'(?:\${1,2})(.*?)(?:\${1,2})'
                for match in re.finditer(formula_pattern, section.content, re.DOTALL):
                    formula_text = match.group(1)
                    # Check if it's inline ($ $) or block ($ $)
                    is_inline = not (match.group(0).startswith('$') and match.group(0).endswith('$'))

                    formula = FormulaData(
                        content=formula_text,
                        is_inline=is_inline,
                        position=match.start(),
                        section_id=section.section_id
                    )
                    document.formulas.append(formula)

    def _extract_formulas_from_text(self, text: str, document: ProcessedDocument) -> None:
        """
            mathematical formulas from plain text.
        """
        if self.config.formula_extraction_method == "pattern":
            # Look for LaTeX style formulas ($...$, $...$)
            formula_pattern = r'(?:\${1,2})(.*?)(?:\${1,2})'
            for match in re.finditer(formula_pattern, text, re.DOTALL):
                formula_text = match.group(1)
                # Check if it's inline ($ $) or block ($ $)
                is_inline = not (match.group(0).startswith('$') and match.group(0).endswith('$'))

                formula = FormulaData(
                    content=formula_text,
                    is_inline=is_inline,
                    position=match.start()
                )
                document.formulas.append(formula)


    def _extract_links_from_text(self, text: str, document: ProcessedDocument) -> None:
        """
        Extract links
        """
        # Look for URLs
        url_pattern = r'(https?://[^\s<>"]+|www\.[^\s<>"]+)'
        for match in re.finditer(url_pattern, text):
            url = match.group(0)
            # Try to find any text that might be associated with the link
            # This is hard to determine from plain text
            link_text = url

            link_data = LinkData(
                text=link_text,
                url=url,
                position=match.start(),
                link_type="external"
            )
            document.links.append(link_data)

        # Look for email addresses
        email_pattern = r'[\w.+-]+@[\w-]+\.[\w.-]+'
        for match in re.finditer(email_pattern, text):
            email = match.group(0)

            link_data = LinkData(
                text=email,
                url=f"mailto:{email}",
                position=match.start(),
                link_type="email"
            )
            document.links.append(link_data)

    def _extract_references_from_document(self, document: ProcessedDocument) -> None:
        """
            Extract references/citations from document content.
        """
        # Iterate through sections
        for section in document.content:
            self._extract_references_from_text(section.content, document, section.section_id)

    def _extract_references_from_text(self, text: str, document: ProcessedDocument, section_id: Optional[str] = None) -> None:
        """
            Extract references from text content.
        """
        # Look for common citation formats

        # Harvard style: (Author, Year)
        harvard_pattern = r'\(([A-Za-z\s]+),\s*(\d{4}[a-z]?)\)'
        for match in re.finditer(harvard_pattern, text):
            author = match.group(1).strip()
            year = match.group(2)

            ref_data = ReferenceData(
                text=match.group(0),
                authors=[author],
                year=int(year[:4]),  # Extract just the numeric part
                position=match.start(),
                section_id=section_id,
                reference_type="citation"
            )
            document.references.append(ref_data)

        # IEEE style: [1], [2], etc.
        ieee_pattern = r'\[(\d+)\]'
        for match in re.finditer(ieee_pattern, text):
            ref_num = match.group(1)

            ref_data = ReferenceData(
                text=match.group(0),
                citation_key=ref_num,
                position=match.start(),
                section_id=section_id,
                reference_type="numbered"
            )
            document.references.append(ref_data)

        # Look for DOIs
        doi_pattern = r'(doi:|DOI:)?\s*(10\.\d{4,}(?:\.\d+)*\/\S+)'
        for match in re.finditer(doi_pattern, text):
            doi = match.group(2).strip()

            ref_data = ReferenceData(
                text=match.group(0),
                doi=doi,
                position=match.start(),
                section_id=section_id,
                reference_type="doi"
            )
            document.references.append(ref_data)

        # Look for references section
        ref_section_pattern = r'(?:References|Bibliography|Works Cited)(?:\s*\n)+((?:(?:[^\n]+\n)+\s*)+)'
        ref_match = re.search(ref_section_pattern, text, re.IGNORECASE)
        if ref_match:
            ref_text = ref_match.group(1)
            # Split by newlines and consider each line or paragraph as a reference
            references = []
            current_ref = []

            for line in ref_text.split('\n'):
                if not line.strip():
                    # Empty line might separate references
                    if current_ref:
                        references.append('\n'.join(current_ref))
                        current_ref = []
                # New reference often starts with number or author name
                elif re.match(r'^\s*\[\d+\]|\d+\.|\w+,', line):
                    if current_ref:
                        references.append('\n'.join(current_ref))
                        current_ref = []
                    current_ref.append(line)
                else:
                    current_ref.append(line)

            # Add the last reference
            if current_ref:
                references.append('\n'.join(current_ref))

            # Create reference objects
            for i, ref in enumerate(references):
                ref_data = ReferenceData(
                    text=ref.strip(),
                    position=ref_match.start() + ref_text.find(ref),
                    section_id=section_id,
                    reference_type="bibliography",
                    metadata={"index": i+1}
                )
                document.references.append(ref_data)

    def _extract_code_blocks_from_document(self, document: ProcessedDocument) -> None:
        """
            Extract code blocks from document content
        """
        # Iterate through sections
        for section in document.content:
            self._extract_code_blocks_from_text(section.content, document)

    def _clean_header_footer(self, text: str, page_num: int, total_pages: int) -> str:
        """
            Clean headers and footers from page text
        """
        if not self.config.clean_header_footer:
            return text

        lines = text.split('\n')
        if len(lines) <= 2:  # Too short to have headers/footers
            return text

        # Basic heuristic for headers/footers:
        # 1. Headers are usually in the first few lines
        # 2. Footers are usually in the last few lines
        # 3. Headers/footers often contain page numbers or document title
        # 4. They often repeat across many pages

        # For this simple version, just remove top and bottom lines if they're short
        header_limit = max(1, int(len(lines) * self.config.max_header_footer_fraction))
        footer_limit = max(1, int(len(lines) * self.config.max_header_footer_fraction))

        # Check for page numbers or short lines
        clean_lines = lines.copy()

        # Check header
        for i in range(min(header_limit, len(lines))):
            line = lines[i].strip()
            # Skip if empty
            if not line:
                continue

            # Check if it's likely a header
            if (len(line) < 100 and 
                (str(page_num + 1) in line or  # Page number
                re.search(r'chapter|section|document|page', line, re.IGNORECASE))):  # Common header terms
                clean_lines[i] = ""

        # Check footer
        for i in range(max(0, len(lines) - footer_limit), len(lines)):
            if i < 0 or i >= len(lines):
                continue

            line = lines[i].strip()
            # Skip if empty
            if not line:
                continue

            # Check if it's likely a footer
            if (len(line) < 100 and 
                (str(page_num + 1) in line or  # Page number
                re.search(r'page|copyright|\d+\s*of\s*\d+', line, re.IGNORECASE))):  # Common footer terms
                clean_lines[i] = ""

        return '\n'.join(clean_lines)

    def _post_process(self, document: ProcessedDocument) -> None:
        """
        method for finalizing document structure and clean extracted content.

        This establishes parent-child relationships between sections,
        cleans text, and links special content to containing sections.
        """
        # Clean up and normalize content
        for section in document.content:
            # Remove excessive whitespace
            section.content = clean_text(section.content)
            # section.content = re.sub(r'\n{3,}', '\n\n', section.content)
            section.content = section.content.strip()

            # Ensure section has non-empty title
            if not section.title.strip():
                # Extract first line or first few words as title
                first_line = section.content.split('\n', 1)[0]
                section.title = (first_line[:50] + '...') if len(first_line) > 50 else first_line

        # Link tables, images, formulas to sections
        self._link_elements_to_sections(document)

        # Set parent-child relationships for sections based on level
        self._set_section_hierarchy(document)

        # Set section paths
        self._set_section_paths(document)

        # Detect languages if configured
        if self.config.detect_languages:
            self._detect_languages(document)

        # Apply custom post-processors if any
        for processor_name, processor_func in self.config.custom_processors.items():
            try:
                processor_func(document)
            except Exception as e:
                logger.error(f"Error in custom processor {processor_name}: {e}")

    def _link_elements_to_sections(self, document: ProcessedDocument) -> None:
        """
            Link tables, images, formulas, code blocks to their respective sections.
        """
        # Skip if no sections
        if not document.content:
            return

        # Check if page information is available for sections
        has_page_info = all(s.page_start is not None for s in document.content)

        # Function to find section by page number
        def find_section_by_page(page_num):
            for section in document.content:
                if (section.page_start <= page_num and 
                    (section.page_end is None or section.page_end >= page_num)):
                    return section.section_id
            return document.content[0].section_id

        # Link tables to sections
        for table in document.tables:
            if not table.section_id:
                if has_page_info and table.page_num is not None:
                    # Find section by page number
                    table.section_id = find_section_by_page(table.page_num)
                else:
                    # No page info, assign to first section
                    table.section_id = document.content[0].section_id

        # Link images to sections
        for image in document.images:
            if not image.section_id:
                if has_page_info and image.page_num is not None:
                    image.section_id = find_section_by_page(image.page_num)
                else:
                    image.section_id = document.content[0].section_id

        # Link formulas to sections
        for formula in document.formulas:
            if not formula.section_id:
                if has_page_info and formula.page_num is not None:
                    formula.section_id = find_section_by_page(formula.page_num)
                else:
                    formula.section_id = document.content[0].section_id

        # Link code blocks to sections
        for code in document.code_blocks:
            if not code.section_id:
                if has_page_info and code.page_num is not None:
                    code.section_id = find_section_by_page(code.page_num)
                else:
                    code.section_id = document.content[0].section_id

        # Link references and links
        for ref in document.references:
            if not ref.section_id:
                if has_page_info and ref.page_num is not None:
                    ref.section_id = find_section_by_page(ref.page_num)

        for link in document.links:
            if not link.section_id:
                if has_page_info and link.page_num is not None:
                    link.section_id = find_section_by_page(link.page_num)

    def _set_section_hierarchy(self, document: ProcessedDocument) -> None:
        """
            Set parent-child relationships for sections based on heading levels.
        """
        if not document.content:
            return

        # Stack to keep track of potential parent sections at each level
        level_to_section = {}

        for section in document.content:
            # Find the most recently seen section with a lower level
            parent_level = max([l for l in level_to_section.keys() if l < section.level], default=0)

            # Set parent if one exists
            if parent_level > 0:
                section.parent_id = level_to_section[parent_level].section_id

            # Update the level map with this section
            level_to_section[section.level] = section

            # Remove any levels higher than current from the map 
            # (as they are siblings or children)
            for level in list(level_to_section.keys()):
                if level > section.level:
                    del level_to_section[level]

    def _set_section_paths(self, document: ProcessedDocument) -> None:
        """
            Set hierarchical paths for each section
        """
        # Create a section id lookup
        section_map = {section.section_id: section for section in document.content}

        # Create path for each section
        for section in document.content:
            path = [section.title]
            current_id = section.parent_id

            # Traverse up the parent chain
            while current_id and current_id in section_map:
                parent = section_map[current_id]
                path.insert(0, parent.title)
                current_id = parent.parent_id

            # Set the path
            section.path = path

    def _detect_languages(self, document: ProcessedDocument) -> None:
        """
            Detect languages used in document sections.
        """
        try:
            import langid
        except ImportError:
            logger.warning("langid library not found. Language detection disabled.")
            return

        # Detect language for each section
        for section in document.content:
            if len(section.content) >= 20:  # Need enough text for reliable detection
                try:
                    lang, confidence = langid.classify(section.content)
                    section.metadata['language'] = lang
                    section.metadata['language_confidence'] = confidence
                except:
                    pass

        # Set document-level language
        text_size = {section.metadata.get('language', ''): len(section.content) 
                    for section in document.content 
                    if 'language' in section.metadata}

        if text_size:
            # Set main language as the one with most content
            main_language = max(text_size.items(), key=lambda x: x[1])[0]
            if main_language:
                document.metadata['primary_language'] = main_language

                # Check if multilingual
                languages = set(text_size.keys())
                if len(languages) > 1:
                    document.metadata['is_multilingual'] = True
                    document.metadata['languages'] = list(languages)
                else:
                    document.metadata['is_multilingual'] = False
